from __future__ import annotations

import hashlib
import io
import os
import struct
import zipfile
import zlib
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from PIL import Image
from pypdf import PdfReader, PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

from app.core.config import get_settings


@pytest.fixture(autouse=True)
def private_root(tmp_path, monkeypatch):
    monkeypatch.setattr(get_settings(), "AI_CHAT_STORAGE_ROOT", str(tmp_path / "private"))
    monkeypatch.setattr(get_settings(), "AI_CHAT_MAX_UPLOAD_BYTES", 4 * 1024 * 1024)
    monkeypatch.setattr(get_settings(), "AI_CHAT_MAX_ATTACHMENT_CHARS", 60_000)


def _service():
    from app.ai_chat import file_service

    return file_service


def _docx_bytes() -> bytes:
    document = Document()
    document.add_paragraph("客户需求：自然黑色")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "长度"
    table.cell(0, 1).text = "数量"
    table.cell(1, 0).text = "18 inch"
    table.cell(1, 1).text = "20"
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _xlsx_bytes() -> bytes:
    workbook = Workbook()
    visible = workbook.active
    visible.title = "报价"
    visible.append(["款号", "数量", "合计"])
    visible.append(["LS-01", 4, "=1+3"])
    hidden = workbook.create_sheet("内部底价")
    hidden["A1"] = "不得泄露"
    hidden.sheet_state = "hidden"
    raw = io.BytesIO()
    workbook.save(raw)
    workbook.close()

    source = zipfile.ZipFile(io.BytesIO(raw.getvalue()))
    output = io.BytesIO()
    with source, zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as target:
        for info in source.infolist():
            data = source.read(info.filename)
            if info.filename == "xl/worksheets/sheet1.xml":
                data = data.replace(b"<f>1+3</f><v></v>", b"<f>1+3</f><v>4</v>")
            target.writestr(info, data)
    return output.getvalue()


def _pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text = "第一张：客户目标"
    table = slide.shapes.add_table(2, 1, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "交期"
    table.cell(1, 0).text = "30 天"
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _pdf_bytes(text: str | None = "Copyable PDF text", *, password: str | None = None) -> bytes:
    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)
    if text is not None:
        font = DictionaryObject(
            {
                NameObject("/Type"): NameObject("/Font"),
                NameObject("/Subtype"): NameObject("/Type1"),
                NameObject("/BaseFont"): NameObject("/Helvetica"),
            }
        )
        page[NameObject("/Resources")] = DictionaryObject(
            {NameObject("/Font"): DictionaryObject({NameObject("/F1"): writer._add_object(font)})}
        )
        stream = DecodedStreamObject()
        escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        stream.set_data(f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET".encode("latin-1"))
        page[NameObject("/Contents")] = writer._add_object(stream)
    if password:
        writer.encrypt(password)
    output = io.BytesIO()
    writer.write(output)
    return output.getvalue()


def _image_bytes(fmt: str, *, exif_orientation: int | None = None) -> bytes:
    image = Image.new("RGB", (2, 3), (200, 10, 20))
    image.info["comment"] = b"private metadata"
    output = io.BytesIO()
    kwargs = {}
    if exif_orientation:
        exif = Image.Exif()
        exif[274] = exif_orientation
        kwargs["exif"] = exif
    image.save(output, format=fmt, **kwargs)
    image.close()
    return output.getvalue()


def _oversized_png_header(width: int, height: int) -> bytes:
    def chunk(kind: bytes, data: bytes) -> bytes:
        body = kind + data
        return struct.pack(">I", len(data)) + body + struct.pack(">I", zlib.crc32(body))

    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    return b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr) + chunk(b"IEND", b"")


@pytest.mark.parametrize(
    ("filename", "mime", "content", "expected"),
    [
        ("notes.txt", "text/plain", b"\xef\xbb\xbfhello\x00 world", "hello world"),
        ("需求.md", "text/markdown", "中文需求\x00说明".encode("gb18030"), "中文需求说明"),
    ],
)
def test_text_and_markdown_decode_bom_legacy_encoding_and_strip_nul(filename, mime, content, expected):
    result = _service().normalize_and_store(filename, mime, content)

    assert result.extracted_text == expected
    assert _service().read_private_file(result.storage_path) == content


def test_docx_extracts_real_paragraphs_and_table_cells():
    result = _service().normalize_and_store(
        "需求.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        _docx_bytes(),
    )

    assert "客户需求：自然黑色" in result.extracted_text
    assert "长度 | 数量" in result.extracted_text
    assert "18 inch | 20" in result.extracted_text


def test_xlsx_uses_cached_formulas_and_visible_nonempty_worksheets_only():
    result = _service().normalize_and_store(
        "报价.xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        _xlsx_bytes(),
    )

    assert "[工作表: 报价]" in result.extracted_text
    assert "LS-01 | 4 | 4" in result.extracted_text
    assert "内部底价" not in result.extracted_text
    assert "不得泄露" not in result.extracted_text


def test_pptx_extracts_slide_text_and_table_text():
    result = _service().normalize_and_store(
        "提案.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        _pptx_bytes(),
    )

    assert "[幻灯片 1]" in result.extracted_text
    assert "第一张：客户目标" in result.extracted_text
    assert "交期" in result.extracted_text and "30 天" in result.extracted_text


def test_pdf_extracts_real_copyable_text():
    content = _pdf_bytes()
    assert "Copyable PDF text" in PdfReader(io.BytesIO(content)).pages[0].extract_text()

    result = _service().normalize_and_store("brief.pdf", "application/pdf", content)

    assert result.extracted_text == "Copyable PDF text"


@pytest.mark.parametrize("content", [_pdf_bytes(password="secret"), _pdf_bytes(None)])
def test_pdf_rejects_encrypted_and_textless_documents_with_actionable_message(content):
    with pytest.raises(_service().FileValidationError, match="可复制文本"):
        _service().normalize_and_store("scan.pdf", "application/pdf", content)


@pytest.mark.parametrize(
    ("fmt", "suffix", "declared", "expected_mime"),
    [
        ("JPEG", ".jpg", "image/jpeg", "image/jpeg"),
        ("PNG", ".png", "image/png", "image/png"),
        ("WEBP", ".webp", "image/webp", "image/webp"),
    ],
)
def test_images_are_magic_validated_normalized_and_metadata_free(fmt, suffix, declared, expected_mime):
    original = _image_bytes(fmt, exif_orientation=6 if fmt == "JPEG" else None)
    result = _service().normalize_and_store(f"portrait{suffix}", declared, original)
    normalized = _service().read_private_file(result.storage_path)

    assert result.mime_type == expected_mime
    assert (result.width, result.height) == ((3, 2) if fmt == "JPEG" else (2, 3))
    assert result.sha256 == hashlib.sha256(normalized).hexdigest()
    with Image.open(io.BytesIO(normalized)) as image:
        image.load()
        assert image.getexif() == {}
        assert "icc_profile" not in image.info
        assert "comment" not in image.info


@pytest.mark.parametrize(
    ("filename", "mime", "content"),
    [
        ("fake.png", "image/png", _image_bytes("JPEG")),
        ("fake.jpg", "image/png", _image_bytes("JPEG")),
        ("legacy.doc", "application/msword", b"legacy"),
        ("notes.markdown", "text/markdown", b"unsupported alias"),
        ("bad.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", b"PK bad"),
        ("wrong.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", _xlsx_bytes()),
    ],
)
def test_spoofed_legacy_corrupt_and_mismatched_formats_are_rejected(filename, mime, content):
    with pytest.raises(_service().FileValidationError):
        _service().normalize_and_store(filename, mime, content)


def test_image_over_60mp_is_rejected_before_decode():
    with pytest.raises(_service().FileValidationError, match="60MP"):
        _service().normalize_and_store(
            "huge.png", "image/png", _oversized_png_header(10_000, 6_001)
        )


def test_upload_over_four_mib_is_rejected_without_writing():
    with pytest.raises(_service().FileValidationError, match="4MiB"):
        _service().normalize_and_store("large.txt", "text/plain", b"x" * (4 * 1024 * 1024 + 1))
    assert not (Path(get_settings().AI_CHAT_STORAGE_ROOT) / "documents").exists()


def test_extraction_truncates_at_configured_limit_with_exact_marker(monkeypatch):
    monkeypatch.setattr(get_settings(), "AI_CHAT_MAX_ATTACHMENT_CHARS", 20)
    marker = "[内容已按系统上限截断]"

    result = _service().normalize_and_store("long.txt", "text/plain", ("内容" * 30).encode())

    assert result.truncated is True
    assert result.extracted_text.endswith(f"\n{marker}")
    assert len(result.extracted_text) <= 20 + 1 + len(marker)


def test_truncated_extraction_closes_resource_backed_iterator(monkeypatch):
    service = _service()

    class TrackedParts:
        closed = False

        def __iter__(self):
            return self

        def __next__(self):
            return "x" * 100

        def close(self):
            self.closed = True

    parts = TrackedParts()
    monkeypatch.setattr(service, "_extract_xlsx", lambda _: parts)
    monkeypatch.setattr(get_settings(), "AI_CHAT_MAX_ATTACHMENT_CHARS", 10)

    service.normalize_and_store(
        "resource.xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        _xlsx_bytes(),
    )

    assert parts.closed is True


def test_storage_uses_private_uuid_relative_path_and_load_delete_helpers():
    result = _service().normalize_and_store("客户名单.txt", None, b"Alice")
    path = Path(result.storage_path)

    assert path.parts[0] == "documents"
    assert path.name != "客户名单.txt"
    assert len(path.stem) == 32 and int(path.stem, 16) >= 0
    assert not path.is_absolute() and "uploads" not in path.parts
    assert _service().resolve_private_path(result.storage_path).read_bytes() == b"Alice"
    _service().delete_private_file(result.storage_path)
    assert not _service().resolve_private_path(result.storage_path).exists()


@pytest.mark.parametrize("relative", ["", ".", "../secret.txt", "/tmp/secret", r"C:\secret.txt"])
def test_resolve_rejects_empty_root_absolute_drive_and_parent_escape(relative):
    with pytest.raises(_service().FileStorageError):
        _service().resolve_private_path(relative)


def test_existing_symlink_path_is_rejected(tmp_path):
    root = Path(get_settings().AI_CHAT_STORAGE_ROOT)
    root.mkdir()
    outside = tmp_path / "outside"
    outside.mkdir()
    link = root / "documents"
    try:
        link.symlink_to(outside, target_is_directory=True)
    except OSError:
        pytest.skip("当前平台不允许创建目录符号链接")

    with pytest.raises(_service().FileStorageError, match="reparse"):
        _service().resolve_private_path("documents/file.txt")


def test_atomic_write_cleans_temporary_file_when_replace_fails(monkeypatch):
    service = _service()
    monkeypatch.setattr(service.os, "replace", lambda *_: (_ for _ in ()).throw(OSError("disk error")))

    with pytest.raises(service.FileStorageError):
        service.normalize_and_store("notes.txt", "text/plain", b"content")
    root = Path(get_settings().AI_CHAT_STORAGE_ROOT)
    assert list(root.rglob("*.tmp")) == [] if root.exists() else True


def test_configured_storage_root_must_be_absolute_for_current_platform(monkeypatch):
    configured = r"D:\WORKSOURCE\ai-chat" if os.name == "posix" else "relative/storage"
    monkeypatch.setattr(get_settings(), "AI_CHAT_STORAGE_ROOT", configured)

    with pytest.raises(_service().FileStorageError, match="绝对路径"):
        _service().resolve_private_path("documents/file.txt")
