"""Validation, extraction, normalization, and private storage for AI chat files."""

from __future__ import annotations

import hashlib
import io
import logging
import os
import stat
import threading
import warnings
import zipfile
from dataclasses import dataclass
from pathlib import Path, PurePosixPath, PureWindowsPath
from uuid import uuid4

from docx import Document
from openpyxl import load_workbook
from PIL import Image, ImageOps, UnidentifiedImageError
from pypdf import PdfReader

from app.core.config import get_settings


MAX_IMAGE_PIXELS = 60_000_000
MAX_OOXML_UNCOMPRESSED_BYTES = 64 * 1024 * 1024
TRUNCATION_MARKER = "[内容已按系统上限截断]"
_STORAGE_LOCK = threading.RLock()
logger = logging.getLogger("commission")

_IMAGE_FORMATS = {
    ".jpg": ("JPEG", "image/jpeg"),
    ".jpeg": ("JPEG", "image/jpeg"),
    ".png": ("PNG", "image/png"),
    ".webp": ("WEBP", "image/webp"),
}
_DOCUMENT_MIMES = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
_TEXT_SUFFIXES = {".txt", ".md", ".markdown"}
_TEXT_MIMES = {"text/plain", "text/markdown", "application/octet-stream"}
_OOXML_PREFIX = {".docx": "word/", ".xlsx": "xl/", ".pptx": "ppt/"}


class FileValidationError(ValueError):
    """The uploaded bytes do not satisfy the accepted attachment contract."""


class FileStorageError(RuntimeError):
    """Private storage is unavailable or a path violates its boundary."""


@dataclass(frozen=True, slots=True)
class StoredAttachment:
    storage_path: str
    original_name: str
    mime_type: str
    file_size: int
    attachment_type: str
    extracted_text: str | None
    truncated: bool
    width: int | None
    height: int | None
    sha256: str


@dataclass(frozen=True, slots=True)
class _PreparedAttachment:
    content: bytes
    suffix: str
    mime_type: str
    attachment_type: str
    extracted_text: str | None = None
    truncated: bool = False
    width: int | None = None
    height: int | None = None


def _normalized_mime(value: str | None) -> str:
    return (value or "").split(";", 1)[0].strip().lower()


def _bounded_text(parts, limit: int) -> tuple[str, bool]:
    kept: list[str] = []
    length = 0
    truncated = False
    iterator = iter(parts)
    try:
        for part in iterator:
            cleaned = str(part).replace("\x00", "").strip()
            if not cleaned:
                continue
            separator = 1 if kept else 0
            available = limit - length - separator
            if available <= 0:
                truncated = True
                break
            if len(cleaned) > available:
                kept.append(cleaned[:available])
                length += separator + available
                truncated = True
                break
            kept.append(cleaned)
            length += separator + len(cleaned)
    finally:
        close = getattr(iterator, "close", None)
        if close is not None:
            close()
    text = "\n".join(kept)
    if truncated:
        text += TRUNCATION_MARKER
    return text, truncated


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8-sig", "gb18030"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            pass
    raise FileValidationError("文本文件编码不受支持，请另存为 UTF-8 后重试")


def _validate_ooxml(content: bytes, suffix: str) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            infos = archive.infolist()
            names = {item.filename for item in infos}
            if "[Content_Types].xml" not in names or not any(
                name.startswith(_OOXML_PREFIX[suffix]) for name in names
            ):
                raise FileValidationError("文件扩展名与实际 Office 格式不匹配")
            if any(item.flag_bits & 1 for item in infos):
                raise FileValidationError("不支持加密的 Office 文件")
            if sum(item.file_size for item in infos) > MAX_OOXML_UNCOMPRESSED_BYTES:
                raise FileValidationError("Office 文件解压后内容过大")
            if archive.testzip() is not None:
                raise FileValidationError("Office 文件已损坏，请重新导出后上传")
    except FileValidationError:
        raise
    except (zipfile.BadZipFile, OSError, ValueError, RuntimeError):
        raise FileValidationError("Office 文件已损坏或格式不正确") from None


def _extract_docx(content: bytes):
    stream = io.BytesIO(content)
    try:
        document = Document(stream)
    finally:
        stream.close()
    for paragraph in document.paragraphs:
        yield paragraph.text
    for table in document.tables:
        for row in table.rows:
            yield " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())


def _extract_xlsx(content: bytes):
    workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    try:
        for worksheet in workbook.worksheets:
            if worksheet.sheet_state != "visible":
                continue
            yield f"[工作表: {worksheet.title}]"
            for row in worksheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value not in (None, "")]
                if values:
                    yield " | ".join(values)
    finally:
        workbook.close()


def _extract_pptx(content: bytes):
    try:
        from pptx import Presentation
    except ImportError:
        raise FileValidationError("PPTX 解析组件未安装，请联系管理员") from None
    stream = io.BytesIO(content)
    try:
        presentation = Presentation(stream)
    finally:
        stream.close()
    for index, slide in enumerate(presentation.slides, 1):
        yield f"[幻灯片 {index}]"
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    yield paragraph.text
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    yield " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())


def _extract_pdf(content: bytes):
    stream = io.BytesIO(content)
    try:
        reader = PdfReader(stream, strict=True)
        if reader.is_encrypted:
            raise FileValidationError("PDF 已加密，无法提取可复制文本，请上传未加密版本")
        for page in reader.pages:
            yield page.extract_text() or ""
    except FileValidationError:
        raise
    except Exception:
        raise FileValidationError("PDF 已损坏或格式不正确") from None
    finally:
        stream.close()


def _prepare_document(suffix: str, declared_mime: str | None, content: bytes) -> _PreparedAttachment:
    mime = _normalized_mime(declared_mime)
    if suffix in _TEXT_SUFFIXES:
        if mime and mime not in _TEXT_MIMES:
            raise FileValidationError("文件 MIME 类型与扩展名不匹配")
        parts = [_decode_text(content)]
        actual_mime = "text/markdown" if suffix in {".md", ".markdown"} else "text/plain"
    else:
        expected_mime = _DOCUMENT_MIMES[suffix]
        if mime and mime != expected_mime:
            raise FileValidationError("文件 MIME 类型与扩展名不匹配")
        actual_mime = expected_mime
        if suffix == ".pdf":
            if not content.startswith(b"%PDF-"):
                raise FileValidationError("文件扩展名与实际 PDF 格式不匹配")
            parts = _extract_pdf(content)
        else:
            _validate_ooxml(content, suffix)
            extractors = {".docx": _extract_docx, ".xlsx": _extract_xlsx, ".pptx": _extract_pptx}
            parts = extractors[suffix](content)
    try:
        text, truncated = _bounded_text(parts, get_settings().AI_CHAT_MAX_ATTACHMENT_CHARS)
    except FileValidationError:
        raise
    except Exception:
        raise FileValidationError("文档已损坏或无法读取，请重新导出后上传") from None
    if not text:
        raise FileValidationError("未提取到可复制文本；如为扫描件，请先进行 OCR 后重试")
    return _PreparedAttachment(content, suffix, actual_mime, "document", text, truncated)


def _magic_image_format(content: bytes) -> str | None:
    if content.startswith(b"\xff\xd8\xff"):
        return "JPEG"
    if content.startswith(b"\x89PNG\r\n\x1a\n"):
        return "PNG"
    if len(content) >= 12 and content[:4] == b"RIFF" and content[8:12] == b"WEBP":
        return "WEBP"
    return None


def _clean_image_mode(image: Image.Image, fmt: str) -> Image.Image:
    if fmt == "JPEG":
        if image.mode in {"RGBA", "LA"} or (image.mode == "P" and "transparency" in image.info):
            rgba = image.convert("RGBA")
            background = Image.new("RGB", rgba.size, "white")
            background.paste(rgba, mask=rgba.getchannel("A"))
            rgba.close()
            return background
        return image.convert("RGB")
    if image.mode == "P":
        return image.convert("RGBA" if "transparency" in image.info else "RGB")
    if image.mode not in {"RGB", "RGBA", "L", "LA"}:
        return image.convert("RGBA" if "A" in image.getbands() else "RGB")
    return image.copy()


def _encode_image(image: Image.Image, fmt: str) -> bytes:
    output = io.BytesIO()
    if fmt == "JPEG":
        image.save(output, "JPEG", quality=90, optimize=True, exif=b"")
    elif fmt == "PNG":
        image.save(output, "PNG", optimize=True)
    else:
        image.save(output, "WEBP", quality=90, method=6, exif=b"", icc_profile=b"")
    return output.getvalue()


def _prepare_image(suffix: str, declared_mime: str | None, content: bytes) -> _PreparedAttachment:
    expected_format, expected_mime = _IMAGE_FORMATS[suffix]
    mime = _normalized_mime(declared_mime)
    if _magic_image_format(content) != expected_format or (mime and mime != expected_mime):
        raise FileValidationError("图片扩展名、MIME 或实际格式不匹配")
    try:
        with warnings.catch_warnings():
            warnings.simplefilter("error", Image.DecompressionBombWarning)
            with Image.open(io.BytesIO(content)) as probe:
                if probe.format != expected_format or probe.width * probe.height > MAX_IMAGE_PIXELS:
                    raise FileValidationError("图片分辨率不能超过 60MP")
                probe.verify()
            with Image.open(io.BytesIO(content)) as source:
                source.load()
                oriented = ImageOps.exif_transpose(source)
                try:
                    cleaned = _clean_image_mode(oriented, expected_format)
                finally:
                    if oriented is not source:
                        oriented.close()
                cleaned.info.clear()
                try:
                    normalized = _encode_image(cleaned, expected_format)
                    width, height = cleaned.size
                finally:
                    cleaned.close()
    except FileValidationError:
        raise
    except (Image.DecompressionBombError, Image.DecompressionBombWarning):
        raise FileValidationError("图片分辨率不能超过 60MP") from None
    except (UnidentifiedImageError, OSError, SyntaxError, ValueError):
        raise FileValidationError("上传内容不是有效图片或图片已损坏") from None
    normalized_suffix = ".jpg" if expected_format == "JPEG" else suffix
    return _PreparedAttachment(
        normalized, normalized_suffix, expected_mime, "image", width=width, height=height
    )


def _storage_root() -> Path:
    configured = str(get_settings().AI_CHAT_STORAGE_ROOT).strip()
    root = Path(configured)
    if not configured or not root.is_absolute():
        raise FileStorageError("私有存储根目录必须是当前系统的绝对路径")
    return Path(os.path.abspath(root))


def _is_reparse_point(path: Path) -> bool:
    try:
        if path.is_symlink() or (hasattr(path, "is_junction") and path.is_junction()):
            return True
        attributes = getattr(path.lstat(), "st_file_attributes", 0)
    except OSError:
        raise FileStorageError("无法验证私有存储路径") from None
    return bool(attributes & getattr(stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0))


def _reject_existing_reparse_points(path: Path) -> None:
    current = Path(path.anchor)
    if os.path.lexists(current) and _is_reparse_point(current):
        raise FileStorageError("私有存储路径包含 reparse point")
    for part in path.parts[1:]:
        current /= part
        if os.path.lexists(current) and _is_reparse_point(current):
            raise FileStorageError("私有存储路径包含 reparse point")


def _resolve_unlocked(relative_path: str) -> Path:
    root = _storage_root()
    _reject_existing_reparse_points(root)
    if not isinstance(relative_path, str) or not relative_path.strip():
        raise FileStorageError("非法私有文件路径")
    windows_path = PureWindowsPath(relative_path)
    posix_path = PurePosixPath(relative_path)
    candidate = Path(relative_path)
    if (
        windows_path.drive
        or windows_path.root
        or posix_path.is_absolute()
        or candidate.is_absolute()
        or any(part == ".." for part in windows_path.parts)
        or any(part == ".." for part in posix_path.parts)
    ):
        raise FileStorageError("非法私有文件路径")
    try:
        target = Path(os.path.abspath(root / candidate))
    except (OSError, ValueError):
        raise FileStorageError("非法私有文件路径") from None
    try:
        inside = os.path.commonpath((str(root), str(target))) == str(root)
    except ValueError:
        inside = False
    if target == root or not inside:
        raise FileStorageError("非法私有文件路径")
    _reject_existing_reparse_points(target)
    return target


def resolve_private_path(relative_path: str) -> Path:
    with _STORAGE_LOCK:
        return _resolve_unlocked(relative_path)


def _cleanup_best_effort(path: Path) -> None:
    try:
        path.unlink(missing_ok=True)
    except Exception as exc:
        message = f"[ai-chat] temporary cleanup failed {path.name}: {exc}"
        logger.warning(message)
        print(message, flush=True)


def _write_atomic(target: Path, content: bytes) -> None:
    temporary = target.with_name(f".{uuid4().hex}.tmp")
    try:
        target.parent.mkdir(parents=True, exist_ok=True)
        _reject_existing_reparse_points(target.parent)
        with temporary.open("xb") as handle:
            handle.write(content)
            handle.flush()
            os.fsync(handle.fileno())
        os.replace(temporary, target)
    except FileStorageError:
        raise
    except OSError as exc:
        message = f"[ai-chat] private attachment write failed {target.name}: {exc}"
        logger.warning(message)
        print(message, flush=True)
        raise FileStorageError("私有文件保存失败，请重试") from exc
    finally:
        _cleanup_best_effort(temporary)


def normalize_and_store(
    filename: str, declared_mime: str | None, content: bytes
) -> StoredAttachment:
    if not isinstance(content, bytes) or not content:
        raise FileValidationError("文件内容为空")
    max_bytes = get_settings().AI_CHAT_MAX_UPLOAD_BYTES
    if len(content) > max_bytes:
        raise FileValidationError(f"单个附件不能超过 {max_bytes // (1024 * 1024)}MiB")
    original_name = PurePosixPath(PureWindowsPath(filename or "").name).name
    suffix = Path(original_name).suffix.lower()
    if suffix in _IMAGE_FORMATS:
        prepared = _prepare_image(suffix, declared_mime, content)
    elif suffix in _TEXT_SUFFIXES or suffix in _DOCUMENT_MIMES:
        prepared = _prepare_document(suffix, declared_mime, content)
    else:
        raise FileValidationError("不支持的文件格式；请上传图片、PDF、DOCX、XLSX、PPTX 或文本")
    category = "images" if prepared.attachment_type == "image" else "documents"
    relative = (Path(category) / f"{uuid4().hex}{prepared.suffix}").as_posix()
    with _STORAGE_LOCK:
        target = _resolve_unlocked(relative)
        _write_atomic(target, prepared.content)
    return StoredAttachment(
        storage_path=relative,
        original_name=original_name,
        mime_type=prepared.mime_type,
        file_size=len(prepared.content),
        attachment_type=prepared.attachment_type,
        extracted_text=prepared.extracted_text,
        truncated=prepared.truncated,
        width=prepared.width,
        height=prepared.height,
        sha256=hashlib.sha256(prepared.content).hexdigest(),
    )


def read_private_file(relative_path: str) -> bytes:
    with _STORAGE_LOCK:
        target = _resolve_unlocked(relative_path)
        try:
            return target.read_bytes()
        except OSError as exc:
            raise FileStorageError("无法读取私有文件") from exc


def delete_private_file(relative_path: str) -> None:
    with _STORAGE_LOCK:
        target = _resolve_unlocked(relative_path)
        try:
            if target.exists() and not target.is_file():
                raise FileStorageError("非法私有文件路径")
            target.unlink(missing_ok=True)
        except FileStorageError:
            raise
        except OSError as exc:
            message = f"[ai-chat] private attachment delete failed {target.name}: {exc}"
            logger.warning(message)
            print(message, flush=True)
            raise FileStorageError("无法删除私有文件") from exc
